"""Build a rich, visually-designed 10-slide thesis defense PPTX.

Native shapes for all diagrams (no imported images): architecture, funnel,
custom bar chart, KPI dashboard, MITRE grid, system stack. Slide transitions
added via XML. Speaker notes preserved in the notes panel.
"""
import json
from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

HERE = Path(__file__).parent
SLIDES_JSON = HERE / "_pptx_slides.json"
OUT_PPTX = HERE / "Aparare_IDS_Kubernetes.pptx"


def C(h):
    return RGBColor((h >> 16) & 0xFF, (h >> 8) & 0xFF, h & 0xFF)


TEAL = C(0x1F8A78); TEAL_LITE = C(0xD7EFEA); TEAL_DARK = C(0x155F52)
ORANGE = C(0xD9822B); ORANGE_LITE = C(0xFCE8D2); ORANGE_DARK = C(0x9C5C16)
PURPLE = C(0x6F4FA3); PURPLE_LITE = C(0xE9E1F4); PURPLE_DARK = C(0x4A3470)
RED = C(0xC83737); RED_LITE = C(0xFDECEC)
GREEN = C(0x2E8B57); GREEN_LITE = C(0xE6F4EB)
GRAY_DARK = C(0x333A40); GRAY_MID = C(0x5B6670)
GRAY_LITE = C(0xEEF1F4); GRAY_LINE = C(0xD0D5DA)
WHITE = C(0xFFFFFF)

SW_IN = 13.333
SH_IN = 7.5


def add_shape(slide, kind, x, y, w, h, fill=None, line=None, line_w=1.0, dashed=False):
    shp = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(line_w)
        if dashed:
            ln = shp.line._get_or_add_ln()
            for existing in ln.findall(qn('a:prstDash')):
                ln.remove(existing)
            pr = etree.SubElement(ln, qn('a:prstDash'))
            pr.set('val', 'dash')
    try:
        shp.shadow.inherit = False
    except Exception:
        pass
    return shp


def add_text(shape, lines, size=14, bold=False, color=GRAY_DARK,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE):
    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(0.05); tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.03); tf.margin_bottom = Inches(0.03)
    if isinstance(lines, str):
        lines = [lines]
    for i, ln in enumerate(lines):
        if isinstance(ln, str):
            text, fmt = ln, {}
        else:
            text, fmt = ln
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = fmt.get('align', align)
        run = p.add_run()
        run.text = text
        run.font.size = Pt(fmt.get('size', size))
        run.font.bold = fmt.get('bold', bold)
        run.font.color.rgb = fmt.get('color', color)
        run.font.name = "Calibri"
    return tf


def add_textbox(slide, x, y, w, h, lines, **kw):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    add_text(box, lines, **kw)
    return box


def down_arrow(slide, x, y, w, h, color):
    return add_shape(slide, MSO_SHAPE.DOWN_ARROW, x, y, w, h, fill=color)


def kpi_card(slide, x, y, w, h, big, label, sub=None, accent=PURPLE, bg=None):
    bg = bg or PURPLE_LITE
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h, fill=bg, line=accent, line_w=1.5)
    add_shape(slide, MSO_SHAPE.RECTANGLE, x, y, w, 0.08, fill=accent)
    add_textbox(slide, x, y + 0.18, w, h * 0.45, big,
                size=32, bold=True, color=accent, align=PP_ALIGN.CENTER)
    add_textbox(slide, x, y + h * 0.6, w, h * 0.25, label,
                size=11, bold=True, color=GRAY_DARK, align=PP_ALIGN.CENTER)
    if sub:
        add_textbox(slide, x + 0.05, y + h * 0.82, w - 0.1, h * 0.18, sub,
                    size=9, color=GRAY_MID, align=PP_ALIGN.CENTER)


def title_band(slide, accent, title, subtitle=None):
    add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, 0.25, SH_IN, fill=accent)
    add_textbox(slide, 0.55, 0.22, 12.5, 0.65, title,
                size=24, bold=True, color=GRAY_DARK, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)
    if subtitle:
        add_textbox(slide, 0.55, 0.9, 12.5, 0.38, subtitle,
                    size=12, color=accent, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)
    add_shape(slide, MSO_SHAPE.RECTANGLE, 0.55, 1.33, 12.2, 0.03, fill=accent)


def footer(slide, page, total, accent):
    add_textbox(slide, 0.55, 7.08, 9.5, 0.32,
                "Apărare teză master · Iulia-Andreea Grigore · iunie 2026",
                size=10, color=GRAY_MID, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)
    add_textbox(slide, 11.6, 7.08, 1.6, 0.32, f"{page} / {total}",
                size=10, bold=True, color=accent, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)


# ---------------- SLIDE 1: TITLE / HERO ----------------
def build_slide_1(slide, data, accent):
    add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, SW_IN, SH_IN, fill=PURPLE)
    for cx, cy, r in [(11.0, 0.5, 2.6), (12.5, 5.8, 1.9), (1.0, 6.7, 1.3), (0.4, 1.2, 0.9)]:
        add_shape(slide, MSO_SHAPE.OVAL, cx - r, cy - r, 2 * r, 2 * r, fill=PURPLE_DARK)
    add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 3.0, SW_IN, 0.05, fill=ORANGE)
    add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 4.55, SW_IN, 0.05, fill=TEAL)

    add_textbox(slide, 0.5, 3.1, 12.333, 1.4,
                "Sistem IDS multi-componentă\npentru Kubernetes",
                size=42, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_textbox(slide, 0.5, 4.65, 12.333, 0.5,
                "Apărare în adâncime pe trafic, audit API și runtime",
                size=18, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_textbox(slide, 0.5, 5.5, 12.333, 0.45,
                "Iulia-Andreea Grigore · sesiunea iunie 2026",
                size=15, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_textbox(slide, 0.5, 5.95, 12.333, 0.4,
                "R3 arhitectură + Flow   ·   R4 Audit + Correlator + deploy",
                size=12, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    badges = ["NIST SP 800-190", "NSA/CISA Hardening v1.2", "MITRE ATT&CK Containers"]
    bw = 2.5; gap = 0.3
    total = bw * len(badges) + gap * (len(badges) - 1)
    sx = (SW_IN - total) / 2
    for i, b in enumerate(badges):
        bx = sx + i * (bw + gap)
        shp = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, bx, 6.7, bw, 0.5,
                        fill=PURPLE_DARK, line=WHITE, line_w=1.0)
        add_text(shp, b, size=11, bold=True, color=WHITE)


# ---------------- SLIDE 2: PROBLEM ----------------
def build_slide_2(slide, data, accent):
    title_band(slide, accent, data["title"], data.get("subtitle"))
    surfaces = [
        ("Rețea L3/L4", "Componenta Flow",
         "DDoS, scanări, anomalii volumetrice",
         "RBAC, secrete, intenții API",
         TEAL, TEAL_LITE),
        ("Audit API", "Componenta Audit",
         "kube-apiserver, RBAC, exec, secrete",
         "pachete, syscalls la kernel",
         ORANGE, ORANGE_LITE),
        ("Syscalls (kernel)", "Componenta Syscall",
         "container escape, file tamper",
         "intenția semantică la nivel API",
         GRAY_MID, GRAY_LITE),
    ]
    cw = 3.95; gap = 0.25
    total = cw * 3 + gap * 2
    sx = (SW_IN - total) / 2
    cy = 1.8; ch = 3.7
    for i, (head, sub, vede, nu, ac, lite) in enumerate(surfaces):
        x = sx + i * (cw + gap)
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, cy, cw, ch,
                  fill=WHITE, line=ac, line_w=1.5)
        add_shape(slide, MSO_SHAPE.RECTANGLE, x, cy, cw, 0.5, fill=ac)
        add_textbox(slide, x, cy + 0.05, cw, 0.4, head,
                    size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_shape(slide, MSO_SHAPE.OVAL, x + cw / 2 - 0.5, cy + 0.7, 1.0, 1.0,
                  fill=lite, line=ac, line_w=1.5)
        add_textbox(slide, x + cw / 2 - 0.5, cy + 0.7, 1.0, 1.0,
                    str(i + 1), size=30, bold=True, color=ac,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, x, cy + 1.85, cw, 0.35, sub,
                    size=12, bold=True, color=ac, align=PP_ALIGN.CENTER)
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x + 0.2, cy + 2.3, cw - 0.4, 0.55,
                  fill=GREEN_LITE, line=GREEN, line_w=0.75)
        add_textbox(slide, x + 0.2, cy + 2.3, cw - 0.4, 0.55,
                    [("✓ vede", {"size": 9, "bold": True, "color": GREEN}),
                     (vede, {"size": 10, "color": GRAY_DARK})],
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x + 0.2, cy + 2.95, cw - 0.4, 0.55,
                  fill=RED_LITE, line=RED, line_w=0.75)
        add_textbox(slide, x + 0.2, cy + 2.95, cw - 0.4, 0.55,
                    [("✗ nu vede", {"size": 9, "bold": True, "color": RED}),
                     (nu, {"size": 10, "color": GRAY_DARK})],
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 2.0, 5.85, 9.333, 0.85,
              fill=PURPLE, line=PURPLE_DARK, line_w=1.5)
    add_textbox(slide, 2.0, 5.85, 9.333, 0.85,
                [("→ niciun sensor unic nu acoperă cele 3 suprafețe", {"size": 16, "bold": True, "color": WHITE}),
                 ("răspuns arhitectural: defense-in-depth · Arp 2022 (USENIX Sec) · NIST SP 800-190", {"size": 10, "color": WHITE})],
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    footer(slide, 2, 10, accent)


# ---------------- SLIDE 3: ARCHITECTURE ----------------
def build_slide_3(slide, data, accent):
    title_band(slide, accent, data["title"], data.get("subtitle"))
    comps = [
        ("Componenta Flow", "XGBoost + Autoencoder", "trafic L3/L4 (R3)", TEAL, TEAL_LITE, False),
        ("Componenta Audit", "Transformer LogBERT", "kube-apiserver (R4)", ORANGE, ORANGE_LITE, False),
        ("Componenta Syscall", "Falco / eBPF", "future work", GRAY_MID, GRAY_LITE, True),
    ]
    cw = 3.5; gap = 0.65
    total = cw * 3 + gap * 2
    sx = (SW_IN - total) / 2
    cy = 1.6; ch = 1.5
    for i, (h, sub, det, ac, lite, dashed) in enumerate(comps):
        x = sx + i * (cw + gap)
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, cy, cw, ch,
                  fill=lite, line=ac, line_w=2.0, dashed=dashed)
        add_textbox(slide, x, cy + 0.12, cw, 0.4, h,
                    size=16, bold=True, color=ac, align=PP_ALIGN.CENTER)
        add_textbox(slide, x, cy + 0.58, cw, 0.4, sub,
                    size=12, color=GRAY_DARK, align=PP_ALIGN.CENTER)
        add_textbox(slide, x, cy + 0.98, cw, 0.4, det,
                    size=10, color=GRAY_MID, align=PP_ALIGN.CENTER)
        down_arrow(slide, x + cw / 2 - 0.18, cy + ch + 0.05, 0.36, 0.55, ac)
    cor_w = 9.5; cor_x = (SW_IN - cor_w) / 2; cor_y = 3.85; cor_h = 1.2
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, cor_x, cor_y, cor_w, cor_h,
              fill=PURPLE, line=PURPLE_DARK, line_w=1.5)
    add_textbox(slide, cor_x, cor_y + 0.1, cor_w, 0.4, "Alert Correlator",
                size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, cor_x, cor_y + 0.55, cor_w, 0.6,
                "threshold → Platt calibration → corelare temporală → MITRE chains → severity",
                size=12, color=WHITE, align=PP_ALIGN.CENTER)
    down_arrow(slide, SW_IN / 2 - 0.18, cor_y + cor_h + 0.1, 0.36, 0.5, PURPLE)
    outs = [("Incidente MITRE", PURPLE, "5 / 3.330 alerte"),
            ("Grafana dashboards", TEAL, "metrici · loguri"),
            ("Email (MailHog)", ORANGE, "alerting nativ")]
    ow = 3.4; ogap = 0.4
    ot = ow * 3 + ogap * 2; osx = (SW_IN - ot) / 2; oy = 5.85
    for i, (t, col, sub) in enumerate(outs):
        x = osx + i * (ow + ogap)
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, oy, ow, 0.85,
                  fill=WHITE, line=col, line_w=1.5)
        add_textbox(slide, x, oy + 0.1, ow, 0.4, t,
                    size=13, bold=True, color=col, align=PP_ALIGN.CENTER)
        add_textbox(slide, x, oy + 0.5, ow, 0.32, sub,
                    size=10, color=GRAY_MID, align=PP_ALIGN.CENTER)
    footer(slide, 3, 10, accent)


# ---------------- SLIDE 4: FLOW DATASETS + HYBRID ----------------
def build_slide_4(slide, data, accent):
    title_band(slide, accent, data["title"], data.get("subtitle"))
    datasets = [
        ("BCCC Cloud DDoS 2024", "Shafi et al. · MDPI",
         "~700k flows · 317 feat · 17 atk TCP", "set principal antrenare",
         TEAL, TEAL_LITE, True),
        ("ITU 2023", "Sever & Dogan",
         "K8s-tagged · CVE-uri reale", "holdout structural (LHO)",
         GRAY_MID, GRAY_LITE, False),
        ("CSE-CIC-IDS 2018", "Canadian Inst.",
         "multi-attack · 77 feat", "novel attack types",
         GRAY_MID, GRAY_LITE, False),
    ]
    cw = 3.85; gap = 0.35
    total = cw * 3 + gap * 2
    sx = (SW_IN - total) / 2
    cy = 1.55; ch = 1.45
    for i, (n, src, sz, role, ac, lite, primary) in enumerate(datasets):
        x = sx + i * (cw + gap)
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, cy, cw, ch,
                  fill=lite, line=ac, line_w=2.5 if primary else 1.0)
        if primary:
            add_shape(slide, MSO_SHAPE.RECTANGLE, x, cy, 0.15, ch, fill=ac)
        add_textbox(slide, x + 0.2, cy + 0.08, cw - 0.3, 0.4, n,
                    size=14, bold=True, color=ac, align=PP_ALIGN.LEFT)
        add_textbox(slide, x + 0.2, cy + 0.5, cw - 0.3, 0.3, src,
                    size=10, color=GRAY_MID, align=PP_ALIGN.LEFT)
        add_textbox(slide, x + 0.2, cy + 0.8, cw - 0.3, 0.3, sz,
                    size=11, bold=True, color=GRAY_DARK, align=PP_ALIGN.LEFT)
        add_textbox(slide, x + 0.2, cy + 1.08, cw - 0.3, 0.3, role,
                    size=10, color=ac, align=PP_ALIGN.LEFT)
    py = 3.15; ph = 0.45
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 1.8, py, 9.733, ph,
              fill=GRAY_LITE, line=GRAY_LINE, line_w=0.75)
    add_textbox(slide, 1.8, py, 9.733, ph,
                "Preprocesare unică · drop IP/port/timestamp anti signature coupling · IANA proto · StandardScaler (AE)",
                size=11, bold=True, color=GRAY_DARK, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    down_arrow(slide, 4.3, 3.68, 0.36, 0.4, TEAL)
    down_arrow(slide, 8.6, 3.68, 0.36, 0.4, ORANGE)
    my = 4.15; mh = 1.5
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 1.8, my, 4.4, mh,
              fill=TEAL, line=TEAL_DARK, line_w=1.5)
    add_textbox(slide, 1.8, my + 0.1, 4.4, 0.4, "XGBoost (supervised)",
                size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, 1.8, my + 0.55, 4.4, 0.4, "Grinsztajn 2022 (NeurIPS)",
                size=10, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, 1.8, my + 0.95, 4.4, 0.4, "« prinde DDoS cunoscut »",
                size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 7.133, my, 4.4, mh,
              fill=ORANGE, line=ORANGE_DARK, line_w=1.5)
    add_textbox(slide, 7.133, my + 0.1, 4.4, 0.4, "Autoencoder (unsupervised)",
                size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, 7.133, my + 0.55, 4.4, 0.4,
                "benign-only · Mirsky 2018 · Kitsune",
                size=10, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, 7.133, my + 0.95, 4.4, 0.4,
                "« prinde devieri de la normal »",
                size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    down_arrow(slide, 4.0, my + mh + 0.05, 0.36, 0.35, TEAL)
    down_arrow(slide, 9.333, my + mh + 0.05, 0.36, 0.35, ORANGE)
    fy = 6.13; fh = 0.7
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 3.0, fy, 7.333, fh,
              fill=PURPLE, line=PURPLE_DARK, line_w=1.5)
    add_textbox(slide, 3.0, fy, 7.333, fh,
                [("Fuziune Platt-calibrată: 0.7 · p_xgb + 0.3 · p_ae", {"size": 14, "bold": True, "color": WHITE})],
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    footer(slide, 4, 10, accent)


# ---------------- SLIDE 5: 3 PATTERNS BAR CHART ----------------
def build_slide_5(slide, data, accent):
    title_band(slide, accent, data["title"], data.get("subtitle"))
    bars = [
        ("BCCC\nheld-out Tuesday", 0.775, "bias temporal", GREEN, "robust"),
        ("CSE-CIC-IDS 2018\ncross-attack", 0.286, "bias spațial", ORANGE, "parțial"),
        ("ITU\nleave-heavy-hitter-out", 0.048, "bias sampling", RED, "catastrofal"),
    ]
    chart_x = 0.7; chart_y = 1.6; chart_w = 7.5; chart_h = 5.0
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, chart_x, chart_y, chart_w, chart_h,
              fill=WHITE, line=GRAY_LINE, line_w=0.5)
    add_textbox(slide, chart_x, chart_y + 0.1, chart_w, 0.35,
                "recall@FPR=1% pe trei protocoale de holdout",
                size=12, bold=True, color=GRAY_DARK, align=PP_ALIGN.CENTER)
    baseline_y = chart_y + chart_h - 1.05
    bar_max_h = baseline_y - (chart_y + 0.55)
    add_shape(slide, MSO_SHAPE.RECTANGLE, chart_x + 0.6, baseline_y,
              chart_w - 1.0, 0.02, fill=GRAY_MID)
    for v, lbl in [(1.0, baseline_y - bar_max_h),
                   (0.5, baseline_y - bar_max_h * 0.5),
                   (0.0, baseline_y)]:
        add_textbox(slide, chart_x + 0.05, lbl - 0.13, 0.5, 0.3,
                    f"{v:.1f}", size=9, color=GRAY_MID, align=PP_ALIGN.RIGHT)
        if v in (0.5, 1.0):
            add_shape(slide, MSO_SHAPE.RECTANGLE, chart_x + 0.6, lbl,
                      chart_w - 1.0, 0.005, fill=GRAY_LINE)
    bar_w = 1.5; bar_gap = 0.65
    bars_total = bar_w * 3 + bar_gap * 2
    bar_start_x = chart_x + 0.85 + (chart_w - 1.0 - bars_total) / 2
    for i, (label, val, bias, col, status) in enumerate(bars):
        bx = bar_start_x + i * (bar_w + bar_gap)
        bh = val * bar_max_h
        by = baseline_y - bh
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, bx, by, bar_w, bh,
                  fill=col, line=col)
        add_textbox(slide, bx - 0.4, by - 0.5, bar_w + 0.8, 0.45,
                    f"{val:.3f}", size=20, bold=True, color=col,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.BOTTOM)
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, bx - 0.05, baseline_y + 0.1,
                  bar_w + 0.1, 0.32, fill=col, line=col)
        add_textbox(slide, bx - 0.05, baseline_y + 0.1, bar_w + 0.1, 0.32,
                    status, size=10, bold=True, color=WHITE,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, bx - 0.4, baseline_y + 0.5, bar_w + 0.8, 0.55,
                    label, size=10, bold=True, color=GRAY_DARK, align=PP_ALIGN.CENTER)
        add_textbox(slide, bx - 0.4, baseline_y + 1.05, bar_w + 0.8, 0.3,
                    bias, size=10, color=col, align=PP_ALIGN.CENTER, bold=True)
    rx = 8.55; rw = 4.4
    kpi_card(slide, rx, 1.6, rw, 1.65, "+2.1pp", "câștig hibrid XGB+AE",
             "recall@FPR=1% 0.782 → 0.803 pe BCCC", accent=PURPLE, bg=PURPLE_LITE)
    kpi_card(slide, rx, 3.4, rw, 1.65, "0.887", "AE singur · ROC-AUC ITU LHO",
             "backstop validat empiric · Sommer-Paxson 2010", accent=ORANGE, bg=ORANGE_LITE)
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, rx, 5.2, rw, 1.4,
              fill=GRAY_LITE, line=GRAY_LINE, line_w=0.75)
    add_textbox(slide, rx + 0.15, 5.3, rw - 0.3, 0.4,
                "Mapare Arp 2022 (USENIX Sec)",
                size=12, bold=True, color=GRAY_DARK, align=PP_ALIGN.LEFT)
    add_textbox(slide, rx + 0.15, 5.7, rw - 0.3, 0.35,
                "§3.4 Spurious Correlations · §3.5 False Causality",
                size=10, color=GRAY_MID, align=PP_ALIGN.LEFT)
    add_textbox(slide, rx + 0.15, 6.05, rw - 0.3, 0.45,
                "eșecurile sunt structurale, nu random",
                size=11, bold=True, color=ORANGE, align=PP_ALIGN.LEFT)
    footer(slide, 5, 10, accent)


# ---------------- SLIDE 6: AUDIT — MITRE GRID + TRANSFORMER ----------------
def build_slide_6(slide, data, accent):
    title_band(slide, accent, data["title"], data.get("subtitle"))
    mitres = [
        ("T1613", "recon", "Discovery"),
        ("T1609", "exec_abuse", "Execution"),
        ("T1078", "rbac_escalation", "Priv Escalation"),
        ("T1552", "secret_access", "Credentials"),
        ("T1528", "sa_token_abuse", "Credentials"),
        ("T1610", "malicious_pod", "Execution"),
    ]
    gx = 0.55; gy = 1.6
    cw = 2.0; ch = 0.85; gap_x = 0.1; gap_y = 0.12
    for idx, (tid, name, cat) in enumerate(mitres):
        row, col = divmod(idx, 3)
        x = gx + col * (cw + gap_x); y = gy + row * (ch + gap_y)
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, cw, ch,
                  fill=ORANGE_LITE, line=ORANGE, line_w=1.5)
        add_textbox(slide, x + 0.1, y + 0.05, cw - 0.2, 0.32, tid,
                    size=12, bold=True, color=ORANGE_DARK, align=PP_ALIGN.LEFT)
        add_textbox(slide, x + 0.1, y + 0.34, cw - 0.2, 0.3, name,
                    size=10, bold=True, color=GRAY_DARK, align=PP_ALIGN.LEFT)
        add_textbox(slide, x + 0.1, y + 0.6, cw - 0.2, 0.25, cat,
                    size=8, color=GRAY_MID, align=PP_ALIGN.LEFT)
    add_textbox(slide, gx, gy + 2 * ch + gap_y + 0.05,
                cw * 3 + gap_x * 2, 0.35,
                "6 scenarii ATT&CK for Containers · 40 runde · 7.396 evenimente · vocab 119",
                size=10, color=GRAY_DARK, align=PP_ALIGN.CENTER, bold=True)
    fy = 4.6
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, gx, fy, 3.05, 1.15,
              fill=C(0xFDF4E8), line=ORANGE, line_w=1.0)
    add_textbox(slide, gx + 0.1, fy + 0.08, 2.85, 0.35,
                "FIX trivialitate (Arp §3.5)",
                size=11, bold=True, color=ORANGE_DARK, align=PP_ALIGN.LEFT)
    add_textbox(slide, gx + 0.1, fy + 0.42, 2.85, 0.7,
                "benign_admin_round cu exec, RBAC, secrete legitime — diferă contextul, nu acțiunea",
                size=9, color=GRAY_DARK, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, gx + 3.15, fy, 3.05, 1.15,
              fill=C(0xFDF4E8), line=ORANGE, line_w=1.0)
    add_textbox(slide, gx + 3.25, fy + 0.08, 2.85, 0.35,
                "FIX script-memo (Arp §3.4)",
                size=11, bold=True, color=ORANGE_DARK, align=PP_ALIGN.LEFT)
    add_textbox(slide, gx + 3.25, fy + 0.42, 2.85, 0.7,
                "randomizare apeluri, ordine, ținte, jitter temporal — nu mai memorează scriptul",
                size=9, color=GRAY_DARK, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)
    rx = 7.0; rw = 6.05
    add_textbox(slide, rx, 1.6, rw, 0.3, "Comparativ 3 modele pe dataset audit",
                size=12, bold=True, color=GRAY_DARK, align=PP_ALIGN.LEFT)
    models = [
        ("XGBoost", 0.86, TEAL_LITE, TEAL, False),
        ("LightGBM", 0.88, TEAL_LITE, TEAL, False),
        ("Transformer (LogBERT-style)", 0.934, ORANGE, ORANGE_DARK, True),
    ]
    my = 1.95
    for i, (mn, f1, fill, line, winner) in enumerate(models):
        y = my + i * 0.78
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, rx, y, 3.5, 0.62,
                  fill=fill, line=line, line_w=2.0 if winner else 1.0)
        add_textbox(slide, rx + 0.15, y, 3.3, 0.62, mn,
                    size=12, bold=winner, color=WHITE if winner else GRAY_DARK,
                    align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)
        track_x = rx + 3.6; track_w = 2.4
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, track_x, y + 0.08, track_w, 0.46,
                  fill=GRAY_LITE, line=GRAY_LINE, line_w=0.5)
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, track_x, y + 0.08, f1 * track_w, 0.46,
                  fill=line, line=line)
        add_textbox(slide, track_x, y + 0.08, track_w, 0.46,
                    f"F1 = {f1:.3f}",
                    size=11, bold=True, color=WHITE,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        if winner:
            add_textbox(slide, rx - 0.4, y, 0.4, 0.62, "★",
                        size=22, bold=True, color=ORANGE,
                        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    kp_y = 4.6; kw = 1.93; gap = 0.13
    kpi_card(slide, rx, kp_y, kw, 1.15, "0.934", "F1", "Transformer",
             accent=ORANGE, bg=ORANGE_LITE)
    kpi_card(slide, rx + (kw + gap), kp_y, kw, 1.15, "0.993", "ROC-AUC", "validation",
             accent=ORANGE, bg=ORANGE_LITE)
    kpi_card(slide, rx + (kw + gap) * 2, kp_y, kw, 1.15, "2.9%", "FPR @ rec 95%", "operating point",
             accent=ORANGE, bg=ORANGE_LITE)
    footer(slide, 6, 10, accent)


# ---------------- SLIDE 7: CORRELATOR — FUNNEL + PIPELINE ----------------
def build_slide_7(slide, data, accent):
    title_band(slide, accent, data["title"], data.get("subtitle"))
    fx = 0.7; fy = 1.65; fw = 4.7
    add_textbox(slide, fx, fy, fw, 0.35,
                "Funnel · 3.330 → 5",
                size=12, bold=True, color=PURPLE, align=PP_ALIGN.LEFT)
    levels = [
        (3330, "alerte brute", GRAY_LITE, GRAY_MID, GRAY_DARK),
        (1024, "post-threshold", TEAL_LITE, TEAL, TEAL_DARK),
        (120, "după corelare 60s", PURPLE_LITE, PURPLE, PURPLE_DARK),
        (5, "incidente MITRE", PURPLE, PURPLE_DARK, WHITE),
    ]
    level_h = 0.85; gap = 0.13
    top_w = fw
    fy0 = fy + 0.4
    for i, (n, lab, fill, line, txt) in enumerate(levels):
        w = top_w * (1 - i * 0.17)
        x = fx + (fw - w) / 2
        y = fy0 + i * (level_h + gap)
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, level_h,
                  fill=fill, line=line, line_w=1.5)
        add_textbox(slide, x, y + 0.05, w, level_h * 0.55,
                    f"{n:,}".replace(",", "."),
                    size=22, bold=True, color=txt, align=PP_ALIGN.CENTER,
                    anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, x, y + level_h * 0.55, w, level_h * 0.4, lab,
                    size=10, color=txt, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    rx = 5.7; rw = 7.4
    add_textbox(slide, rx, 1.65, rw, 0.3, "Pipeline 5 niveluri",
                size=12, bold=True, color=PURPLE, align=PP_ALIGN.LEFT)
    steps = [("threshold", TEAL), ("Platt", TEAL),
             ("corelare 60s", PURPLE), ("MITRE chains", PURPLE),
             ("severity", PURPLE_DARK)]
    sx = rx; sy = 2.0; sw = 1.3; sh = 0.65; sgap = 0.12
    for i, (t, col) in enumerate(steps):
        x = sx + i * (sw + sgap)
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, sy, sw, sh, fill=col, line=col)
        add_textbox(slide, x, sy, sw, sh, t,
                    size=10, bold=True, color=WHITE,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        if i < len(steps) - 1:
            add_shape(slide, MSO_SHAPE.RIGHT_ARROW, x + sw + 0.005,
                      sy + sh / 2 - 0.08, 0.11, 0.16, fill=GRAY_MID)
    cy_ = 2.95
    add_textbox(slide, rx, cy_, rw, 0.3,
                "Lanțuri MITRE detectate (multiplicatori severity)",
                size=12, bold=True, color=PURPLE, align=PP_ALIGN.LEFT)
    chains = [
        ("full_kill_chain", "T1613 → T1609 → T1078 → T1610", "3.0×", RED),
        ("recon_to_escape", "T1613 → T1610", "2.0×", ORANGE),
        ("credential_to_escalation", "T1552/T1528 → T1078", "1.8×", PURPLE),
    ]
    yy = 3.3
    for i, (n, seq, mult, col) in enumerate(chains):
        y = yy + i * 0.7
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, rx, y, rw, 0.6,
                  fill=WHITE, line=col, line_w=1.5)
        add_textbox(slide, rx + 0.15, y + 0.05, 3.0, 0.5, n,
                    size=12, bold=True, color=col,
                    align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, rx + 3.2, y + 0.05, 2.9, 0.5, seq,
                    size=10, color=GRAY_MID,
                    align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, rx + rw - 0.95, y + 0.1,
                  0.85, 0.4, fill=col, line=col)
        add_textbox(slide, rx + rw - 0.95, y + 0.1, 0.85, 0.4, mult,
                    size=12, bold=True, color=WHITE,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    kp_y = 5.65; kw = 2.35; gap_k = 0.12
    kpi_card(slide, rx, kp_y, kw, 1.15, "99.85%", "dedup", "3.330 → 5",
             accent=PURPLE, bg=PURPLE_LITE)
    kpi_card(slide, rx + (kw + gap_k), kp_y, kw, 1.15, "100%", "precision", "incident-level",
             accent=PURPLE, bg=PURPLE_LITE)
    kpi_card(slide, rx + (kw + gap_k) * 2, kp_y, kw, 1.15, "6/6", "tipuri atac", "acoperite",
             accent=PURPLE, bg=PURPLE_LITE)
    footer(slide, 7, 10, accent)


# ---------------- SLIDE 8: SYSTEM ----------------
def build_slide_8(slide, data, accent):
    title_band(slide, accent, data["title"], data.get("subtitle"))
    cx = 0.55; cy = 1.6; cw = 8.3; ch = 5.1
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, cx, cy, cw, ch,
              fill=GRAY_LITE, line=GRAY_MID, line_w=1.5, dashed=True)
    add_textbox(slide, cx + 0.15, cy + 0.05, cw - 0.3, 0.35,
                "Cluster Kubernetes (AKS · kind)",
                size=12, bold=True, color=GRAY_DARK, align=PP_ALIGN.LEFT)
    pw = 3.7; pgap = 0.3; py = cy + 0.55; ph = 1.7
    px1 = cx + 0.3; px2 = px1 + pw + pgap
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, px1, py, pw, ph,
              fill=TEAL_LITE, line=TEAL, line_w=1.5)
    add_shape(slide, MSO_SHAPE.RECTANGLE, px1, py, pw, 0.05, fill=TEAL)
    add_textbox(slide, px1, py + 0.1, pw, 0.4,
                "Flow IDS Pod",
                size=14, bold=True, color=TEAL_DARK, align=PP_ALIGN.CENTER)
    add_textbox(slide, px1, py + 0.5, pw, 0.32,
                "FastAPI · non-root",
                size=10, color=GRAY_MID, align=PP_ALIGN.CENTER)
    add_textbox(slide, px1, py + 0.85, pw, 0.4,
                "/predict   /healthz   /readyz   /metrics",
                size=10, color=GRAY_DARK, align=PP_ALIGN.CENTER)
    add_textbox(slide, px1, py + 1.25, pw, 0.4,
                "XGBoost + Autoencoder · Platt",
                size=11, bold=True, color=TEAL_DARK, align=PP_ALIGN.CENTER)
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, px2, py, pw, ph,
              fill=ORANGE_LITE, line=ORANGE, line_w=1.5)
    add_shape(slide, MSO_SHAPE.RECTANGLE, px2, py, pw, 0.05, fill=ORANGE)
    add_textbox(slide, px2, py + 0.1, pw, 0.4, "Audit IDS Pod",
                size=14, bold=True, color=ORANGE_DARK, align=PP_ALIGN.CENTER)
    add_textbox(slide, px2, py + 0.5, pw, 0.32,
                "FastAPI · non-root",
                size=10, color=GRAY_MID, align=PP_ALIGN.CENTER)
    add_textbox(slide, px2, py + 0.85, pw, 0.4,
                "/predict   /healthz   /readyz   /metrics",
                size=10, color=GRAY_DARK, align=PP_ALIGN.CENTER)
    add_textbox(slide, px2, py + 1.25, pw, 0.4,
                "Transformer LogBERT · 20 tokens",
                size=11, bold=True, color=ORANGE_DARK, align=PP_ALIGN.CENTER)
    dy = py + ph + 0.25; dh = 1.05; dw = pw * 2 + pgap; dx = px1
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, dx, dy, dw, dh,
              fill=PURPLE_LITE, line=PURPLE, line_w=1.5)
    add_textbox(slide, dx, dy + 0.1, dw, 0.4,
                "DaemonSet · audit log streamer (1 pod / nod)",
                size=13, bold=True, color=PURPLE_DARK, align=PP_ALIGN.CENTER)
    add_textbox(slide, dx, dy + 0.55, dw, 0.45,
                "tail hostPath audit.log · ferestre per actor · publish la Audit IDS Pod",
                size=10, color=GRAY_DARK, align=PP_ALIGN.CENTER)
    ny = dy + dh + 0.18
    add_textbox(slide, cx + 0.15, ny, cw - 0.3, 0.35,
                "Namespace · Deployment · Service · ConfigMap · NetworkPolicy · ServiceMonitor",
                size=10, color=GRAY_MID, align=PP_ALIGN.CENTER, bold=True)
    ox = 9.15; ow = 3.95
    add_textbox(slide, ox, 1.6, ow, 0.35, "Observabilitate (Helm)",
                size=13, bold=True, color=PURPLE, align=PP_ALIGN.LEFT)
    stacks = [
        ("Prometheus", "metrici", TEAL, TEAL_LITE),
        ("Loki", "logs", ORANGE, ORANGE_LITE),
        ("Grafana", "dashboards + alerting", PURPLE, PURPLE_LITE),
    ]
    sy = 2.0; sh_s = 0.95; sgap = 0.18
    for i, (n, d, col, lite) in enumerate(stacks):
        y = sy + i * (sh_s + sgap)
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, ox, y, ow, sh_s,
                  fill=lite, line=col, line_w=1.5)
        add_shape(slide, MSO_SHAPE.RECTANGLE, ox, y, 0.12, sh_s, fill=col)
        add_textbox(slide, ox + 0.25, y + 0.12, ow - 0.35, 0.35, n,
                    size=14, bold=True, color=col, align=PP_ALIGN.LEFT)
        add_textbox(slide, ox + 0.25, y + 0.55, ow - 0.35, 0.35, d,
                    size=10, color=GRAY_DARK, align=PP_ALIGN.LEFT)
    ey = sy + 3 * (sh_s + sgap) + 0.05
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, ox, ey, ow, 0.7,
              fill=GRAY_DARK, line=GRAY_DARK)
    add_textbox(slide, ox, ey, ow, 0.7,
                [("✉ Email (MailHog dev)", {"size": 13, "bold": True, "color": WHITE}),
                 ("alerting nativ Grafana · dedup la incident", {"size": 9, "color": WHITE})],
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    footer(slide, 8, 10, accent)


# ---------------- SLIDE 9: CONCLUSIONS DASHBOARD ----------------
def build_slide_9(slide, data, accent):
    title_band(slide, accent, data["title"], data.get("subtitle"))
    kpis = [
        ("3", "pattern-uri Arp 2022",
         "temporal · spațial · sampling — triangulare empirică", TEAL),
        ("+2.1pp", "câștig hibrid XGB+AE",
         "recall@FPR=1% pe BCCC, fără cost FPR", TEAL),
        ("0.934", "F1 Transformer audit",
         "ROC-AUC 0.993 · FPR 2.9% · 6 scenarii MITRE", ORANGE),
        ("99.85%", "dedup Alert Correlator",
         "3.330 alerte → 5 incidente · precision 100%", PURPLE),
    ]
    kx = 0.55; ky = 1.55; kw = 6.15; kh = 2.15; kgap_x = 0.3; kgap_y = 0.3
    for i, (big, label, sub, col) in enumerate(kpis):
        r, c = divmod(i, 2)
        x = kx + c * (kw + kgap_x); y = ky + r * (kh + kgap_y)
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, kw, kh,
                  fill=WHITE, line=col, line_w=2.0)
        add_shape(slide, MSO_SHAPE.RECTANGLE, x, y, 0.2, kh, fill=col)
        add_textbox(slide, x + 0.4, y + 0.2, 3.0, kh - 0.4, big,
                    size=44, bold=True, color=col,
                    align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, x + 3.5, y + 0.3, kw - 3.65, 0.55, label,
                    size=14, bold=True, color=GRAY_DARK,
                    align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, x + 3.5, y + 0.95, kw - 3.65, 1.0, sub,
                    size=11, color=GRAY_MID, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)
    hy = 6.4
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 0.55, hy, 12.2, 0.6,
              fill=GRAY_DARK, line=GRAY_DARK)
    add_textbox(slide, 0.55, hy, 12.2, 0.6,
                "Defense-in-depth = răspunsul corect arhitectural, nu un model mai mare",
                size=15, bold=True, color=WHITE,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    footer(slide, 9, 10, accent)


# ---------------- SLIDE 10: LIMITATIONS + FUTURE + Q&A ----------------
def build_slide_10(slide, data, accent):
    title_band(slide, accent, data["title"], data.get("subtitle"))
    lims = [
        ("Signature coupling", "ITU LHO recall 0.048",
         "Flow ≠ detector universal"),
        ("Class balance ~23%", "dataset audit",
         "nereal pentru producție — recalibrare necesară"),
        ("Falco respins pe AKS", "kernel 5.15-azure",
         "modern_ebpf verificator + kmod 0 alerte"),
    ]
    lx = 0.55; ly = 1.55; lw = 4.1; lh = 1.65; lgap = 0.2
    for i, (t, body, sub) in enumerate(lims):
        x = lx + i * (lw + lgap)
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, ly, lw, lh,
                  fill=RED_LITE, line=RED, line_w=1.5)
        add_shape(slide, MSO_SHAPE.ISOSCELES_TRIANGLE, x + 0.15, ly + 0.2, 0.55, 0.55, fill=RED)
        add_textbox(slide, x + 0.15, ly + 0.32, 0.55, 0.4, "!",
                    size=20, bold=True, color=WHITE,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, x + 0.85, ly + 0.18, lw - 0.95, 0.4, t,
                    size=14, bold=True, color=RED, align=PP_ALIGN.LEFT)
        add_textbox(slide, x + 0.85, ly + 0.58, lw - 0.95, 0.4, body,
                    size=11, color=GRAY_DARK, align=PP_ALIGN.LEFT)
        add_textbox(slide, x + 0.15, ly + 1.1, lw - 0.3, 0.45, sub,
                    size=10, color=GRAY_MID, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)
    fy = 3.5
    add_textbox(slide, 0.55, fy, 12.2, 0.4, "Future work",
                size=14, bold=True, color=PURPLE, align=PP_ALIGN.LEFT)
    add_shape(slide, MSO_SHAPE.RECTANGLE, 1.0, fy + 1.0, 11.3, 0.04, fill=GRAY_MID)
    futs = [
        ("Falco + Syscall", "kernel compatibil"),
        ("Federated learning", "clustere multiple"),
        ("Validare externă", "dataset Audit"),
        ("Latency under load", "AKS multi-node"),
    ]
    nw = len(futs); slot = 10.5 / nw
    for i, (n, sub) in enumerate(futs):
        cx_ = 1.5 + i * slot
        add_shape(slide, MSO_SHAPE.OVAL, cx_ - 0.22, fy + 0.8, 0.44, 0.44,
                  fill=PURPLE, line=WHITE, line_w=2.0)
        add_textbox(slide, cx_ - 1.5, fy + 0.45, 3.0, 0.32, n,
                    size=11, bold=True, color=PURPLE, align=PP_ALIGN.CENTER)
        add_textbox(slide, cx_ - 1.5, fy + 1.3, 3.0, 0.32, sub,
                    size=9, color=GRAY_MID, align=PP_ALIGN.CENTER)
    hy = 5.45
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 0.55, hy, 12.2, 1.4,
              fill=PURPLE, line=PURPLE)
    add_textbox(slide, 0.55, hy + 0.15, 12.2, 0.6, "Q & A",
                size=44, bold=True, color=WHITE,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_textbox(slide, 0.55, hy + 0.85, 12.2, 0.4,
                "Mulțumesc coordonatorului, familiei și comisiei",
                size=14, color=WHITE,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    footer(slide, 10, 10, GRAY_DARK)


BUILDERS = {
    1: build_slide_1, 2: build_slide_2, 3: build_slide_3, 4: build_slide_4, 5: build_slide_5,
    6: build_slide_6, 7: build_slide_7, 8: build_slide_8, 9: build_slide_9, 10: build_slide_10,
}
ACCENTS = {
    1: PURPLE, 2: GRAY_MID, 3: PURPLE, 4: TEAL, 5: TEAL,
    6: ORANGE, 7: PURPLE, 8: PURPLE, 9: GRAY_DARK, 10: GRAY_DARK,
}

NS_P = "http://schemas.openxmlformats.org/presentationml/2006/main"


def add_fade_transition(slide):
    xml = f'<p:transition xmlns:p="{NS_P}" spd="med"><p:fade/></p:transition>'
    slide.element.append(etree.fromstring(xml))


def add_entrance_animation(slide):
    """Fade-in entrance on all shapes, triggered on slide click."""
    shape_ids = [sh.shape_id for sh in slide.shapes if sh.shape_id is not None]
    if not shape_ids:
        return
    timing_xml = f'''<p:timing xmlns:p="{NS_P}" xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
  <p:tnLst>
    <p:par>
      <p:cTn id="1" dur="indefinite" restart="never" nodeType="tmRoot">
        <p:childTnLst>
          <p:seq concurrent="1" nextAc="seek">
            <p:cTn id="2" dur="indefinite" nodeType="mainSeq">
              <p:childTnLst>
'''
    cid = 3
    for sid in shape_ids:
        timing_xml += f'''                <p:par>
                  <p:cTn id="{cid}" fill="hold">
                    <p:stCondLst><p:cond delay="indefinite"/></p:stCondLst>
                    <p:childTnLst>
                      <p:par>
                        <p:cTn id="{cid+1}" fill="hold">
                          <p:stCondLst><p:cond delay="0"/></p:stCondLst>
                          <p:childTnLst>
                            <p:par>
                              <p:cTn id="{cid+2}" presetID="10" presetClass="entr" presetSubtype="0" fill="hold" grpId="0" nodeType="clickEffect">
                                <p:stCondLst><p:cond delay="0"/></p:stCondLst>
                                <p:childTnLst>
                                  <p:set>
                                    <p:cBhvr>
                                      <p:cTn id="{cid+3}" dur="1" fill="hold"><p:stCondLst><p:cond delay="0"/></p:stCondLst></p:cTn>
                                      <p:tgtEl><p:spTgt spid="{sid}"/></p:tgtEl>
                                      <p:attrNameLst><p:attrName>style.visibility</p:attrName></p:attrNameLst>
                                    </p:cBhvr>
                                    <p:to><p:strVal val="visible"/></p:to>
                                  </p:set>
                                  <p:anim calcmode="lin" valueType="num">
                                    <p:cBhvr additive="base">
                                      <p:cTn id="{cid+4}" dur="300" fill="hold"/>
                                      <p:tgtEl><p:spTgt spid="{sid}"/></p:tgtEl>
                                      <p:attrNameLst><p:attrName>style.opacity</p:attrName></p:attrNameLst>
                                    </p:cBhvr>
                                    <p:tavLst>
                                      <p:tav tm="0"><p:val><p:fltVal val="0"/></p:val></p:tav>
                                      <p:tav tm="100000"><p:val><p:fltVal val="1"/></p:val></p:tav>
                                    </p:tavLst>
                                  </p:anim>
                                </p:childTnLst>
                              </p:cTn>
                            </p:par>
                          </p:childTnLst>
                        </p:cTn>
                      </p:par>
                    </p:childTnLst>
                  </p:cTn>
                </p:par>
'''
        cid += 5
    timing_xml += '''              </p:childTnLst>
            </p:cTn>
          </p:seq>
        </p:childTnLst>
      </p:cTn>
    </p:par>
  </p:tnLst>
</p:timing>'''
    slide.element.append(etree.fromstring(timing_xml))


def main():
    slides_data = json.loads(SLIDES_JSON.read_text())
    assert len(slides_data) == 10

    prs = Presentation()
    prs.slide_width = Inches(SW_IN)
    prs.slide_height = Inches(SH_IN)

    blank = prs.slide_layouts[6]

    for i, sd in enumerate(slides_data, 1):
        slide = prs.slides.add_slide(blank)
        accent = ACCENTS[i]
        BUILDERS[i](slide, sd, accent)
        notes = slide.notes_slide.notes_text_frame
        notes.text = ""
        p = notes.paragraphs[0]
        r = p.add_run(); r.text = sd["speakerNotes"]
        r.font.size = Pt(12); r.font.name = "Calibri"
        add_fade_transition(slide)

    prs.save(OUT_PPTX)
    print(f"Wrote {OUT_PPTX} ({OUT_PPTX.stat().st_size/1024:.1f} KB)")


if __name__ == "__main__":
    main()
